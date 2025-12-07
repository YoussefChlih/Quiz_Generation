"""Document Processing Module
Handles extraction of text from various document formats (PDF, PPTX, DOCX, TXT, RTF)
"""

import os
import re
from typing import Optional

# PDF extraction - try multiple libraries for best results
try:
    import pdfplumber
    HAS_PDFPLUMBER = True
except ImportError:
    HAS_PDFPLUMBER = False

import PyPDF2
from pptx import Presentation
from docx import Document
from striprtf.striprtf import rtf_to_text


class DocumentProcessor:
    """Process various document types and extract text content"""
    
    SUPPORTED_EXTENSIONS = {'.pdf', '.pptx', '.ppt', '.docx', '.doc', '.txt', '.rtf'}
    
    def __init__(self):
        self.extractors = {
            '.pdf': self._extract_pdf,
            '.pptx': self._extract_pptx,
            '.ppt': self._extract_pptx,
            '.docx': self._extract_docx,
            '.doc': self._extract_docx,
            '.txt': self._extract_txt,
            '.rtf': self._extract_rtf
        }
    
    def process(self, file_path: str) -> Optional[str]:
        """
        Extract text from a document file
        
        Args:
            file_path: Path to the document file
            
        Returns:
            Extracted text content or None if extraction fails
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(f"Unsupported file format: {ext}")
        
        extractor = self.extractors.get(ext)
        if extractor:
            return extractor(file_path)
        
        return None
    
    def _extract_pdf(self, file_path: str) -> str:
        """Extract text from PDF files using multiple methods for best results"""
        text_content = []
        
        # Method 1: Try pdfplumber first (better for complex PDFs)
        if HAS_PDFPLUMBER:
            try:
                with pdfplumber.open(file_path) as pdf:
                    for page_num, page in enumerate(pdf.pages):
                        page_text = page.extract_text()
                        if page_text and page_text.strip():
                            # Clean the text
                            cleaned_text = self._clean_text(page_text)
                            if cleaned_text:
                                text_content.append(f"--- Page {page_num + 1} ---\n{cleaned_text}")
                        
                        # Also extract tables if present
                        tables = page.extract_tables()
                        for table in tables:
                            if table:
                                table_text = self._table_to_text(table)
                                if table_text:
                                    text_content.append(table_text)
                
                if text_content:
                    return '\n\n'.join(text_content)
            except Exception:
                pass  # Fall back to PyPDF2
        
        # Method 2: Fallback to PyPDF2
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                for page_num, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    if page_text and page_text.strip():
                        cleaned_text = self._clean_text(page_text)
                        if cleaned_text:
                            text_content.append(f"--- Page {page_num + 1} ---\n{cleaned_text}")
                        
        except Exception as e:
            raise Exception(f"Error extracting PDF: {str(e)}")
        
        return '\n\n'.join(text_content)
    
    def _clean_text(self, text: str) -> str:
        """Clean extracted text"""
        if not text:
            return ""
        
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        # Remove special characters that might cause issues
        text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]', '', text)
        # Fix common OCR issues
        text = text.replace('\ufeff', '').replace('\u200b', '')
        # Normalize line breaks
        text = re.sub(r'\n\s*\n', '\n\n', text)
        
        return text.strip()
    
    def _table_to_text(self, table: list) -> str:
        """Convert table data to readable text"""
        if not table:
            return ""
        
        rows = []
        for row in table:
            if row:
                cells = [str(cell).strip() if cell else '' for cell in row]
                if any(cells):  # Only add non-empty rows
                    rows.append(' | '.join(cells))
        
        return '\n'.join(rows) if rows else ""
    
    def _extract_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint files"""
        text_content = []
        
        try:
            prs = Presentation(file_path)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                slide_text.append(f"--- Slide {slide_num} ---")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                    
                    # Extract text from tables
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = [cell.text for cell in row.cells if cell.text]
                            if row_text:
                                slide_text.append(' | '.join(row_text))
                
                text_content.append('\n'.join(slide_text))
                
        except Exception as e:
            raise Exception(f"Error extracting PPTX: {str(e)}")
        
        return '\n\n'.join(text_content)
    
    def _extract_docx(self, file_path: str) -> str:
        """Extract text from Word documents"""
        text_content = []
        
        try:
            doc = Document(file_path)
            
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text_content.append(para.text)
            
            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text for cell in row.cells if cell.text.strip()]
                    if row_text:
                        text_content.append(' | '.join(row_text))
                        
        except Exception as e:
            raise Exception(f"Error extracting DOCX: {str(e)}")
        
        return '\n\n'.join(text_content)
    
    def _extract_txt(self, file_path: str) -> str:
        """Extract text from plain text files"""
        try:
            # Try different encodings
            encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        return file.read()
                except UnicodeDecodeError:
                    continue
            
            raise Exception("Could not decode file with any supported encoding")
            
        except Exception as e:
            raise Exception(f"Error extracting TXT: {str(e)}")
    
    def _extract_rtf(self, file_path: str) -> str:
        """Extract text from RTF files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                rtf_content = file.read()
            
            return rtf_to_text(rtf_content)
            
        except Exception as e:
            raise Exception(f"Error extracting RTF: {str(e)}")


def get_file_info(file_path: str) -> dict:
    """Get information about a document file"""
    if not os.path.exists(file_path):
        return None
    
    stat = os.stat(file_path)
    
    return {
        'name': os.path.basename(file_path),
        'size': stat.st_size,
        'extension': os.path.splitext(file_path)[1].lower(),
        'path': file_path
    }
